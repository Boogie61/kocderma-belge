
# ============================================================
#  kocderma.com — Belge Üretim Sunucusu (FastAPI)
#  Claude içeriği JSON şema olarak üretir; bu sunucu gerçek
#  docx / xlsx / pptx / pdf dosyasını kendi üretir ve indirtir.
#  Hiçbir beta/özel API özelliğine bağlı değildir.
#
#  Ortam değişkeni: ANTHROPIC_API_KEY (Render'da Secret olarak ekle)
# ============================================================
 
import os
import re
import json
import uuid
import time
import base64
 
import requests
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel
 
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.lib import colors
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem, Table, TableStyle,
)
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
 
ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
MODEL = "claude-haiku-4-5-20251001"        # yedek (JSON) yontemi icin
SKILL_MODEL = "claude-sonnet-4-6"          # resmi Skill'ler icin (kalite)
FILES_API = "https://api.anthropic.com/v1/files"
SKILL_BETAS = "code-execution-2025-08-25,skills-2025-10-02,files-api-2025-04-14"
SKILLS_LIST = [{"type": "anthropic", "skill_id": _s, "version": "latest"} for _s in ["pptx", "docx", "xlsx", "pdf"]]
FILES_DIR = os.environ.get("FILES_DIR", "/tmp/genfiles")
os.makedirs(FILES_DIR, exist_ok=True)
 
ALLOWED_ORIGINS = [
    "https://kocderma.com",
    "https://www.kocderma.com",
]
 
# Türkçe karakterler için Unicode font (PDF). Linux'ta DejaVu genelde kuruludur.
PDF_FONT = "Helvetica"
for _fp in [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
]:
    if os.path.exists(_fp):
        try:
            pdfmetrics.registerFont(TTFont("DejaVu", _fp))
            PDF_FONT = "DejaVu"
            break
        except Exception:
            pass
 
app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_methods=["*"],
    allow_headers=["*"],
)
 
SYSTEM = """Sen bir belge üreticisisin. Kullanıcının isteğine göre TEK BİR JSON nesnesi üret.
Yanıtın SADECE JSON olsun: açıklama yazma, kod bloğu işareti (```) kullanma.
 
Şema:
{
  "kind": "docx" | "xlsx" | "pptx" | "pdf",
  "filename": "uygun-ad.uzanti",
  "title": "Belge başlığı",
  "summary": "Kullanıcıya gösterilecek 1 cümlelik Türkçe özet",
  "blocks": [
     {"type":"heading","level":1,"text":"..."},
     {"type":"paragraph","text":"..."},
     {"type":"bullets","items":["...","..."]},
     {"type":"numbered","items":["...","..."]},
     {"type":"table","headers":["..."],"rows":[["...","..."]]}
  ],
  "subtitle": "Sunum kapağı için kısa alt başlık (pptx)",
  "sheets": [ {"name":"Sayfa1","headers":["..."],"rows":[["..."]]} ],
  "slides": [
     {"title":"Bölüm başlığı","layout":"section"},
     {"title":"İçerik slaytı","layout":"content","bullets":["...","..."]},
     {"title":"Karşılaştırma","layout":"two-col","columns":[
        {"heading":"Sütun A","bullets":["...","..."]},
        {"heading":"Sütun B","bullets":["...","..."]}]},
     {"title":"Tablo slaytı","layout":"table","headers":["...","..."],"rows":[["...","..."]]},
     {"title":"Vurgu","layout":"callout","text":"Vurgulanacak tek anahtar cümle"},
     {"title":"Sayılarla","layout":"stats","stats":[
        {"value":"%80","label":"kısa açıklama"},{"value":"3:1","label":"kısa açıklama"}]}
  ]
}
 
Kurallar:
- Kullanıcı format belirttiyse (Word=docx, Excel=xlsx, sunum/PowerPoint=pptx, PDF=pdf) onu kullan; belirtmediyse içeriğe en uygun olanı seç.
- docx ve pdf için "blocks" doldur. xlsx için "sheets". pptx için "title" + "subtitle" + "slides".
- İçeriği Türkçe, dolu ve dermatoloji eğitimine uygun hazırla; önemli İngilizce terimleri parantez içinde ekle.
- Bu bir eğitim aracıdır; tanı koyma, kişiye özel tıbbi tavsiye verme.
 
Sunum (pptx) için ek kurallar:
- "title" konuyu, "subtitle" kısa bir tanımı versin — kapak slaytı bunlardan üretilir.
- "slides" mantıklı bir akış olsun: her ana konu öbeğinden önce bir bölüm ayracı ekle
  ("layout":"section", sadece "title" — bullets YOK).
- ÖNEMLİ — slayt tiplerini ÇEŞİTLENDİR; sürekli madde listesi monoton durur. Şu tipleri uygun
  yerlerde kullan:
  · "content" → 3-6 kısa madde (öz ifadeler, tam paragraf değil)
  · "two-col" → iki kavramı/yaklaşımı yan yana karşılaştırırken (columns + heading)
  · "table"   → sınıflama, dozaj, ayırıcı tanı gibi yapısal veriyi gösterirken
  · "callout" → bir slaytı tek bir anahtar mesaj/klinik inci ile vurgularken (text)
  · "stats"   → prevalans, oran, sağkalım gibi sayısal bilgileri öne çıkarırken (stats)
- Tipik bir sunumda en az 2-3 farklı tip kullan. Slaytı metin duvarına çevirme.
- Kullanıcı belirli bir şey isterse (tablo, karşılaştırma, vurgu, istatistik...) o tipi kullan.
- Kapsamlı sunum istenirse 12-20 slayt üret; en sonda bir "Özet / Anahtar Noktalar" slaytı ekle.
- SADECE geçerli JSON döndür."""
 
 
DOC_SKILL_SYSTEM = (
    "Sen kocderma.com icin bir belge asistanisin. Kullanicinin istedigi belgeyi UYGUN Skill ile "
    "olustur: PowerPoint (pptx), Word (docx), Excel (xlsx) veya PDF. Profesyonel, tasarimli ve dolu "
    "bir cikti uret. Icerik Turkce ve dermatoloji egitimine uygun olsun; onemli Ingilizce terimleri "
    "parantez icinde ekle. Asil icerigi DOSYAYA koy; sohbet yanitin kisa olsun. Bu bir egitim aracidir; "
    "tani koyma, kisiye ozel tibbi tavsiye verme."
)
 
 
class DocReq(BaseModel):
    messages: list
 
 
def call_claude(messages):
    body = {"model": MODEL, "max_tokens": 8000, "system": SYSTEM, "messages": messages}
    last = "Claude API hatasi"
    for attempt in range(4):
        try:
            r = requests.post(
                "https://api.anthropic.com/v1/messages",
                headers={
                    "x-api-key": ANTHROPIC_API_KEY,
                    "anthropic-version": "2023-06-01",
                    "content-type": "application/json",
                },
                json=body,
                timeout=90,
            )
        except Exception as e:
            last = "baglanti/zaman asimi: " + str(e)
            time.sleep(2 * (attempt + 1))
            continue
        if r.status_code == 200:
            data = r.json()
            parts = [b.get("text", "") for b in data.get("content", []) if b.get("type") == "text"]
            return "\n".join(parts)
        if r.status_code in (429, 529, 503):
            last = "Sunucu yogun (" + str(r.status_code) + ")"
            time.sleep(2 * (attempt + 1))
            continue
        try:
            last = r.json().get("error", {}).get("message", "") or ("hata " + str(r.status_code))
        except Exception:
            last = "hata " + str(r.status_code)
        break
    raise RuntimeError(last)
 
 
def parse_spec(text):
    t = (text or "").strip()
    t = re.sub(r"^```[a-zA-Z]*", "", t).strip()
    t = re.sub(r"```$", "", t).strip()
    start, end = t.find("{"), t.rfind("}")
    if start >= 0 and end > start:
        t = t[start:end + 1]
    return json.loads(t)
 
 
def build_docx(spec, path):
    doc = Document()
    if spec.get("title"):
        doc.add_heading(str(spec["title"]), 0)
    for b in spec.get("blocks", []) or []:
        tp = b.get("type")
        if tp == "heading":
            lvl = min(max(int(b.get("level", 1) or 1), 1), 4)
            doc.add_heading(str(b.get("text", "")), lvl)
        elif tp == "paragraph":
            doc.add_paragraph(str(b.get("text", "")))
        elif tp == "bullets":
            for it in b.get("items", []) or []:
                doc.add_paragraph(str(it), style="List Bullet")
        elif tp == "numbered":
            for it in b.get("items", []) or []:
                doc.add_paragraph(str(it), style="List Number")
        elif tp == "table":
            headers = b.get("headers", []) or []
            rows = b.get("rows", []) or []
            ncol = len(headers) or (len(rows[0]) if rows else 0)
            if ncol:
                table = doc.add_table(rows=0, cols=ncol)
                try:
                    table.style = "Table Grid"
                except Exception:
                    pass
                if headers:
                    cells = table.add_row().cells
                    for i, h in enumerate(headers):
                        if i < ncol:
                            cells[i].text = str(h)
                for row in rows:
                    cells = table.add_row().cells
                    for i, c in enumerate(row):
                        if i < ncol:
                            cells[i].text = str(c)
    doc.save(path)
 
 
def build_xlsx(spec, path):
    wb = Workbook()
    sheets = spec.get("sheets")
    if not sheets:
        sheets = [{"name": "Sayfa1", "headers": spec.get("headers", []), "rows": spec.get("rows", [])}]
    first = True
    for sh in sheets:
        ws = wb.active if first else wb.create_sheet()
        first = False
        ws.title = (str(sh.get("name") or "Sayfa"))[:31]
        headers = sh.get("headers", []) or []
        rows = sh.get("rows", []) or []
        if headers:
            ws.append([str(h) for h in headers])
        for row in rows:
            ws.append([str(c) for c in row])
    wb.save(path)
 
 
# ---- pptx klinik-mavi tema ----
PX_NAVY  = RGBColor(0x0E, 0x1A, 0x2B)   # koyu lacivert
PX_BLUE  = RGBColor(0x1F, 0x5F, 0xB3)   # klinik mavi (vurgu)
PX_BLUED = RGBColor(0x14, 0x3F, 0x7A)   # koyu mavi
PX_LIGHT = RGBColor(0xF7, 0xF8, 0xFA)   # açık zemin
PX_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PX_INK   = RGBColor(0x1A, 0x24, 0x33)   # metin
PX_DIM   = RGBColor(0x5A, 0x66, 0x78)   # soluk metin
PX_SOFT  = RGBColor(0xDD, 0xE7, 0xF5)   # açık mavi
PX_TINT  = RGBColor(0xEC, 0xF2, 0xFA)   # çok açık mavi (kart zemini)
PX_FONT  = "Calibri"
PX_SW    = Inches(13.333)               # 16:9 genişlik
PX_SH    = Inches(7.5)                  # 16:9 yükseklik
 
PX_LAYOUTS = ("section", "content", "two-col", "table", "callout", "stats")
 
 
def _px_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color
 
 
def _px_rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp
 
 
def _px_round(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp
 
 
def _px_box(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf
 
 
def _px_run(p, text, size, color, bold=False, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = PX_FONT
    return r
 
 
def _px_bullets(tf, items, size=18, gap=12):
    for i, b in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.16
        _px_run(p, "▪  ", size, PX_BLUE, bold=True)
        _px_run(p, str(b), size, PX_INK)
 
 
def _px_header(slide, s_title, deck_title, idx, total):
    """Açık zeminli içerik slaytı iskeleti (üst çubuk, başlık, altbilgi, sayfa no)."""
    _px_bg(slide, PX_LIGHT)
    _px_rect(slide, 0, 0, PX_SW, Inches(0.16), PX_BLUE)
    tf = _px_box(slide, Inches(0.85), Inches(0.5), Inches(11.6), Inches(1.0))
    _px_run(tf.paragraphs[0], s_title, 28, PX_NAVY, bold=True)
    _px_rect(slide, Inches(0.9), Inches(1.5), Inches(1.4), Inches(0.055), PX_BLUE)
    _px_rect(slide, Inches(0.9), Inches(6.84), Inches(11.5), Inches(0.02), PX_SOFT)
    tf = _px_box(slide, Inches(0.9), Inches(6.9), Inches(9.5), Inches(0.4))
    _px_run(tf.paragraphs[0], deck_title + "  ·  kocderma.com", 10, PX_DIM)
    tf = _px_box(slide, Inches(11.75), Inches(6.86), Inches(1.2), Inches(0.4))
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _px_run(tf.paragraphs[0], "%d / %d" % (idx, total), 10, PX_DIM)
 
 
def build_pptx(spec, path):
    prs = Presentation()
    prs.slide_width = PX_SW
    prs.slide_height = PX_SH
    blank = prs.slide_layouts[6]
 
    title = str(spec.get("title") or "Sunum")
    subtitle = str(spec.get("subtitle") or "")
    slides = spec.get("slides", []) or []
 
    # ---- Kapak slaytı ----
    cover = prs.slides.add_slide(blank)
    _px_bg(cover, PX_BLUE)
    _px_rect(cover, 0, 0, Inches(0.28), PX_SH, PX_BLUED)
    tf = _px_box(cover, Inches(1.0), Inches(2.25), Inches(11.3), Inches(0.5))
    _px_run(tf.paragraphs[0], "BOLOGNIA DERMATOLOJİ  ·  kocderma.com", 14, PX_SOFT, bold=True)
    tf = _px_box(cover, Inches(1.0), Inches(2.75), Inches(11.4), Inches(2.3))
    _px_run(tf.paragraphs[0], title, 44, PX_WHITE, bold=True)
    _px_rect(cover, Inches(1.05), Inches(4.95), Inches(2.2), Inches(0.06), PX_WHITE)
    if subtitle:
        tf = _px_box(cover, Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.4))
        _px_run(tf.paragraphs[0], subtitle, 19, PX_SOFT)
 
    total = len(slides)
    section_no = 0
    for idx0, sl in enumerate(slides):
        idx = idx0 + 1
        s_title = str(sl.get("title", "") or "")
        bullets = sl.get("bullets", []) or []
        layout = str(sl.get("layout", "") or "").lower()
        if layout not in PX_LAYOUTS:
            layout = "content" if bullets else "section"
 
        slide = prs.slides.add_slide(blank)
 
        # ---- Bölüm ayracı ----
        if layout == "section":
            section_no += 1
            _px_bg(slide, PX_NAVY)
            _px_rect(slide, 0, 0, Inches(0.28), PX_SH, PX_BLUE)
            tf = _px_box(slide, Inches(1.0), Inches(2.05), Inches(5), Inches(1.4))
            _px_run(tf.paragraphs[0], "%02d" % section_no, 54, PX_BLUE, bold=True)
            tf = _px_box(slide, Inches(1.0), Inches(3.05), Inches(11.3), Inches(2.2))
            _px_run(tf.paragraphs[0], s_title, 36, PX_WHITE, bold=True)
            _px_rect(slide, Inches(1.05), Inches(4.55), Inches(1.8), Inches(0.06), PX_BLUE)
            tf = _px_box(slide, Inches(11.75), Inches(6.86), Inches(1.2), Inches(0.4))
            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
            _px_run(tf.paragraphs[0], "%d / %d" % (idx, total), 10, PX_SOFT)
            continue
 
        # ---- İçerik tipli slaytlar (ortak iskelet) ----
        _px_header(slide, s_title, title, idx, total)
        cy = Inches(1.95)
 
        # ---- İki sütun / karşılaştırma ----
        if layout == "two-col":
            cols = sl.get("columns", []) or []
            if not cols and bullets:
                half = (len(bullets) + 1) // 2
                cols = [{"bullets": bullets[:half]}, {"bullets": bullets[half:]}]
            xs = [Inches(0.9), Inches(7.0)]
            cw = Inches(5.45)
            for ci, col in enumerate(cols[:2]):
                yy = cy
                heading = str(col.get("heading", "") or "")
                if heading:
                    _px_rect(slide, xs[ci], yy, cw, Inches(0.5), PX_BLUE)
                    htf = _px_box(slide, xs[ci] + Inches(0.15), yy, cw - Inches(0.3),
                                  Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
                    _px_run(htf.paragraphs[0], heading, 15, PX_WHITE, bold=True)
                    yy = yy + Inches(0.7)
                btf = _px_box(slide, xs[ci], yy, cw, Inches(4.4))
                _px_bullets(btf, col.get("bullets", []) or [], size=16, gap=9)
 
        # ---- Tablo ----
        elif layout == "table":
            headers = sl.get("headers", []) or []
            rows = sl.get("rows", []) or []
            data = ([headers] if headers else []) + [list(r) for r in rows]
            data = [r for r in data if r]
            if data:
                ncol = max(len(r) for r in data)
                nrow = len(data)
                th = min(Inches(0.55) * nrow, Inches(4.6))
                gt = slide.shapes.add_table(nrow, ncol, Inches(0.9), cy,
                                            Inches(11.5), th).table
                for ri, row in enumerate(data):
                    is_h = bool(headers) and ri == 0
                    for ci in range(ncol):
                        cell = gt.cell(ri, ci)
                        cell.text = ""
                        val = row[ci] if ci < len(row) else ""
                        _px_run(cell.text_frame.paragraphs[0], str(val), 13,
                                PX_WHITE if is_h else PX_INK, bold=is_h)
                        cell.fill.solid()
                        if is_h:
                            cell.fill.fore_color.rgb = PX_BLUE
                        else:
                            cell.fill.fore_color.rgb = PX_TINT if ri % 2 else PX_WHITE
                        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
 
        # ---- Vurgu kutusu (callout) ----
        elif layout == "callout":
            text = str(sl.get("text", "") or (bullets[0] if bullets else ""))
            _px_round(slide, Inches(1.4), Inches(2.7), Inches(10.5), Inches(2.3), PX_TINT)
            _px_rect(slide, Inches(1.4), Inches(2.78), Inches(0.12), Inches(2.14), PX_BLUE)
            tf = _px_box(slide, Inches(2.0), Inches(2.7), Inches(9.4), Inches(2.3),
                         anchor=MSO_ANCHOR.MIDDLE)
            _px_run(tf.paragraphs[0], text, 23, PX_NAVY, bold=True)
            extra = list(bullets)
            if not sl.get("text") and extra:
                extra = extra[1:]
            if extra:
                etf = _px_box(slide, Inches(1.5), Inches(5.3), Inches(10.3), Inches(1.4))
                _px_bullets(etf, extra, size=15, gap=7)
 
        # ---- İstatistik kartları ----
        elif layout == "stats":
            stats = (sl.get("stats", []) or [])[:4]
            n = len(stats) or 1
            gap = Inches(0.3)
            cardw = (Inches(11.5) - gap * (n - 1)) // n
            for si, st in enumerate(stats):
                x = Inches(0.9) + (cardw + gap) * si
                _px_round(slide, x, Inches(2.55), cardw, Inches(2.6), PX_WHITE, line=PX_SOFT)
                _px_rect(slide, x, Inches(2.55), cardw, Inches(0.12), PX_BLUE)
                vtf = _px_box(slide, x, Inches(3.0), cardw, Inches(1.2),
                              anchor=MSO_ANCHOR.MIDDLE)
                vtf.paragraphs[0].alignment = PP_ALIGN.CENTER
                _px_run(vtf.paragraphs[0], str(st.get("value", "")), 38, PX_BLUE, bold=True)
                ltf = _px_box(slide, x + Inches(0.2), Inches(4.15), cardw - Inches(0.4),
                              Inches(0.9))
                ltf.paragraphs[0].alignment = PP_ALIGN.CENTER
                _px_run(ltf.paragraphs[0], str(st.get("label", "")), 14, PX_DIM)
            if bullets:
                btf = _px_box(slide, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.3))
                _px_bullets(btf, bullets, size=15, gap=7)
 
        # ---- Klasik madde slaytı ----
        else:
            btf = _px_box(slide, Inches(0.9), cy, Inches(11.5), Inches(4.7))
            _px_bullets(btf, bullets, size=18, gap=12)
 
    prs.save(path)
 
 
def build_pdf(spec, path):
    doc = SimpleDocTemplate(path, pagesize=A4, leftMargin=20 * mm, rightMargin=20 * mm,
                            topMargin=18 * mm, bottomMargin=18 * mm)
    styles = getSampleStyleSheet()
    for sn in ["Title", "Heading1", "Heading2", "Heading3", "BodyText"]:
        styles[sn].fontName = PDF_FONT
    story = []
    if spec.get("title"):
        story.append(Paragraph(str(spec["title"]), styles["Title"]))
        story.append(Spacer(1, 8))
    for b in spec.get("blocks", []) or []:
        tp = b.get("type")
        if tp == "heading":
            lvl = min(max(int(b.get("level", 1) or 1), 1), 3)
            story.append(Paragraph(str(b.get("text", "")), styles["Heading" + str(lvl)]))
        elif tp == "paragraph":
            story.append(Paragraph(str(b.get("text", "")), styles["BodyText"]))
            story.append(Spacer(1, 4))
        elif tp in ("bullets", "numbered"):
            items = [ListItem(Paragraph(str(it), styles["BodyText"])) for it in (b.get("items", []) or [])]
            story.append(ListFlowable(items, bulletType="1" if tp == "numbered" else "bullet"))
            story.append(Spacer(1, 4))
        elif tp == "table":
            headers = b.get("headers", []) or []
            rows = b.get("rows", []) or []
            data = []
            if headers:
                data.append([str(h) for h in headers])
            for row in rows:
                data.append([str(c) for c in row])
            if data:
                tbl = Table(data, hAlign="LEFT")
                tbl.setStyle(TableStyle([
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F5FB3")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONTNAME", (0, 0), (-1, -1), PDF_FONT),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 5),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 5),
                ]))
                story.append(tbl)
                story.append(Spacer(1, 6))
    if not story:
        story.append(Paragraph(str(spec.get("summary", "Belge")), styles["BodyText"]))
    doc.build(story)
 
 
def _anthropic_headers(betas):
    h = {"x-api-key": ANTHROPIC_API_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"}
    if betas:
        h["anthropic-beta"] = betas
    return h
 
 
def _collect_file_ids(content):
    ids = []
    for block in content or []:
        if not isinstance(block, dict):
            continue
        if "code_execution_tool_result" in block.get("type", ""):
            c = block.get("content")
            inner = c.get("content") if isinstance(c, dict) else None
            if isinstance(inner, list):
                for o in inner:
                    if isinstance(o, dict) and o.get("file_id"):
                        ids.append(o["file_id"])
    return ids
 
 
def _download_anthropic_file(fid):
    h = _anthropic_headers("files-api-2025-04-14")
    name = "belge"
    try:
        mr = requests.get(FILES_API + "/" + fid, headers=h, timeout=60)
        if mr.status_code == 200:
            name = mr.json().get("filename", name) or name
    except Exception:
        pass
    cr = requests.get(FILES_API + "/" + fid + "/content", headers=h, timeout=180)
    cr.raise_for_status()
    return name, cr.content
 
 
def skills_generate(messages):
    """Anthropic resmi Skill'leri (code execution) ile gercek dosya uretir."""
    container = {"skills": SKILLS_LIST}
    msgs = list(messages)
    data = {}
    for _ in range(8):
        body = {
            "model": SKILL_MODEL,
            "max_tokens": 12000,
            "system": DOC_SKILL_SYSTEM,
            "messages": msgs,
            "tools": [{"type": "code_execution_20250825", "name": "code_execution"}],
            "container": container,
        }
        r = requests.post("https://api.anthropic.com/v1/messages",
                          headers=_anthropic_headers(SKILL_BETAS), json=body, timeout=240)
        if r.status_code != 200:
            try:
                msg = r.json().get("error", {}).get("message", "")
            except Exception:
                msg = (r.text or "")[:200]
            raise RuntimeError("skill api " + str(r.status_code) + ": " + (msg or ""))
        data = r.json()
        cid = (data.get("container") or {}).get("id")
        if cid:
            container = {"id": cid, "skills": SKILLS_LIST}
        if data.get("stop_reason") == "pause_turn":
            msgs = msgs + [{"role": "assistant", "content": data.get("content", [])}]
            continue
        break
    reply = "\n".join(b.get("text", "") for b in data.get("content", []) if b.get("type") == "text").strip()
    files = []
    for fid in _collect_file_ids(data.get("content", [])):
        name, blob = _download_anthropic_file(fid)
        stored = uuid.uuid4().hex + "__" + name
        with open(os.path.join(FILES_DIR, stored), "wb") as f:
            f.write(blob)
        files.append({"id": stored, "name": name})
    return reply, files
 
 
BUILDERS = {"docx": build_docx, "xlsx": build_xlsx, "pptx": build_pptx, "pdf": build_pdf}
 
 
@app.get("/")
def health():
    return {"ok": True, "service": "kocderma-belge", "key": bool(ANTHROPIC_API_KEY)}
 
 
@app.post("/doc")
def make_doc(req: DocReq):
    if not ANTHROPIC_API_KEY:
        return {"reply": "Sunucuda API anahtari ayarli degil.", "files": []}
    if not req.messages:
        return {"reply": "Bos istek.", "files": []}
 
    # 1) Resmi Skill'lerle uret (en iyi kalite)
    skill_err = ""
    try:
        s_reply, s_files = skills_generate(req.messages)
        if s_files:
            return {"reply": s_reply or "Belge hazir.", "files": s_files}
        skill_err = "skill dosya uretmedi"
    except Exception as e:
        skill_err = str(e)
 
    # 2) Yedek: JSON taslagi + python ile uret
    try:
        text = call_claude(req.messages)
        spec = parse_spec(text)
    except Exception as e:
        return {"reply": "Belge olusturulamadi (skill: " + skill_err + " / yedek: " + str(e) + ")", "files": []}
 
    kind = str(spec.get("kind") or "docx").lower()
    if kind not in BUILDERS:
        kind = "docx"
    filename = str(spec.get("filename") or ("belge." + kind))
    filename = re.sub(r"[^A-Za-z0-9_.\-çğıöşüÇĞİÖŞÜ ]", "", filename) or ("belge." + kind)
    if not filename.lower().endswith("." + kind):
        filename = filename + "." + kind
 
    tmp = os.path.join(FILES_DIR, uuid.uuid4().hex + "." + kind)
    try:
        BUILDERS[kind](spec, tmp)
        with open(tmp, "rb") as fh:
            data = fh.read()
    except Exception as e:
        return {"reply": "Dosya uretiminde hata: " + str(e), "files": []}
    finally:
        try:
            if os.path.exists(tmp):
                os.remove(tmp)
        except Exception:
            pass
 
    # Dosyayı yanıtın içinde doğrudan (base64) gönderiyoruz. Render'ın geçici
    # diskine güvenmiyoruz — ikinci bir indirme isteği olmadığı için "dosya yok"
    # hatası ve cold-start/uyku arası kopma riski ortadan kalkar.
    return {
        "reply": str(spec.get("summary") or "Belge hazir."),
        "files": [{
            "name": filename,
            "kind": kind,
            "size": len(data),
            "b64": base64.b64encode(data).decode("ascii"),
        }],
    }
 
 
@app.get("/files/{file_id}")
def get_file(file_id: str):
    if "/" in file_id or "\\" in file_id or ".." in file_id:
        raise HTTPException(status_code=400, detail="gecersiz id")
    path = os.path.join(FILES_DIR, file_id)
    if not os.path.exists(path):
        raise HTTPException(status_code=404, detail="dosya yok")
    name = file_id.split("__", 1)[1] if "__" in file_id else file_id
    return FileResponse(path, filename=name)
    name = file_id.split("__", 1)[1] if "__" in file_id else file_id
    return FileResponse(path, filename=name)
